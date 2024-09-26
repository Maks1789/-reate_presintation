import openpyxl
import os
from pptx import Presentation
from pptx.util import Pt
from docx import Document


class ExcelProcessor:
    def __init__(self, file_path):
        self.file_path = file_path
        self.urls = set()

    def load_data(self):
        wb = openpyxl.load_workbook(self.file_path)
        sheet = wb.active
        self._collect_urls(sheet)
        self._clear_sheet(sheet)
        self._save_unique_urls(sheet, wb)

    def _collect_urls(self, sheet):
        prefixes = ["https:", "tiktok"]
        for row in sheet.iter_rows(min_col=1, max_col=1, values_only=True):
            cell_value = str(row[0])
            if cell_value and cell_value.startswith(tuple(prefixes)):
                self.urls.add(cell_value)

    def _clear_sheet(self, sheet):
        for row in sheet.iter_rows(min_col=1, max_col=1):
            for cell in row:
                cell.value = None

    def _save_unique_urls(self, sheet, workbook):
        for idx, url in enumerate(sorted(self.urls, key=len), start=1):
            sheet.cell(row=idx, column=1, value=url)
        new_file_name = "processed_" + os.path.basename(self.file_path)
        new_file_path = os.path.join(os.path.dirname(self.file_path), new_file_name)
        workbook.save(new_file_path)


class DocxProcessor:
    def __init__(self, file_path):
        self.file_path = file_path
        self.urls = set()

    def load_data(self):
        doc = Document(self.file_path)
        self._collect_urls(doc)

    def _collect_urls(self, doc):
        prefixes = ["https:", "tiktok"]
        for paragraph in doc.paragraphs:
            words = paragraph.text.split()
            for word in words:
                if any(word.startswith(prefix) for prefix in prefixes):
                    self.urls.add(word)


class PowerPointProcessor:
    def __init__(self, file_path):
        self.file_path = file_path
        self.prs = Presentation(file_path)
        self.chunk_size = 15
        self.text_box_index = 5

    def process_urls(self, urls):
        slide_index = 0
        while urls:
            current_chunk = urls[:self.chunk_size]
            urls = urls[self.chunk_size:]

            if any(len(url) > 110 for url in current_chunk):
                self.chunk_size = 10

            if slide_index >= len(self.prs.slides):
                break

            slide = self.prs.slides[slide_index]
            text_frame = slide.shapes[self.text_box_index].text_frame
            text_frame.clear()

            for url in current_chunk:
                self.add_paragraph(text_frame, url)

            slide_index += 1

    def add_paragraph(self, text_frame, text):
        p = text_frame.add_paragraph()
        p.text = text
        p.font.name = 'Montserrat'
        p.font.size = Pt(28)
        p.space_before = Pt(0)
        p.space_after = Pt(0)

        empty_paragraph = text_frame.add_paragraph()
        empty_paragraph.text = ''
        empty_paragraph.font.name = 'Montserrat'
        empty_paragraph.font.size = Pt(28)
        empty_paragraph.space_before = Pt(0)
        empty_paragraph.space_after = Pt(0)

    def save_presentation(self):
        output_pptx_file = "processed_" + os.path.basename(self.file_path)
        output_pptx_path = os.path.join(os.path.dirname(self.file_path), output_pptx_file)
        self.prs.save(output_pptx_path)


if __name__ == "__main__":
    location_files = os.listdir(os.getcwd())
    all_find_file_pptx = []
    all_find_file_excel = []
    all_find_file_docx = []

    # Фільтрація файлів
    def file_filter(location_files):
        try:
            for file in location_files:
                if file.endswith(".pptx"):
                    all_find_file_pptx.append(str(file))
                elif file.endswith(".xlsx"):
                    all_find_file_excel.append(str(file))
                elif file.endswith(".docx"):
                    all_find_file_docx.append(str(file))
        except Exception as e:
            print(f"Помилка при фільтрації файлів: {e}")

    file_filter(location_files)

    # Об'єднання всіх посилань
    all_urls = set()

    # Обробка файлів Excel
    for excel_file in all_find_file_excel:
        excel_processor = ExcelProcessor(excel_file)
        excel_processor.load_data()
        all_urls.update(excel_processor.urls)

    # Обробка файлів Docx
    for docx_file in all_find_file_docx:
        docx_processor = DocxProcessor(docx_file)
        docx_processor.load_data()
        all_urls.update(docx_processor.urls)

    # Обробка файлу PowerPoint
    if all_find_file_pptx:
        pptx_file = all_find_file_pptx[0]
        ppt_processor = PowerPointProcessor(pptx_file)
        sorted_urls = sorted(all_urls, key=len)
        ppt_processor.process_urls(sorted_urls)
        ppt_processor.save_presentation()
